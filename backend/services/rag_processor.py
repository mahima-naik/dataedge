import io
from loguru import logger

# Conditional imports for optional dependencies
try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

def extract_text_from_file(content: bytes, filename: str) -> str:
    """Extracts text from various file formats for RAG injection."""
    ext = filename.split('.')[-1].lower()
    
    if ext == 'pdf':
        if not PdfReader:
            return "[Error: pypdf not installed]"
        reader = PdfReader(io.BytesIO(content))
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    
    elif ext == 'docx':
        if not Document:
            return "[Error: python-docx not installed]"
        doc = Document(io.BytesIO(content))
        text = "\n".join([para.text for para in doc.paragraphs])
        return text.strip()
    
    elif ext == 'pptx':
        if not Presentation:
            return "[Error: python-pptx not installed]"
        prs = Presentation(io.BytesIO(content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    
    elif ext == 'txt':
        return content.decode('utf-8', errors='ignore').strip()
    
    else:
        return f"[Unsupported file format: {ext}]"
